# *******Install pptx python library with yor current python version using command- (pip install python-pptx) ***********
# ********First install sample_font_file.ttf in your system*************

import collections
import collections.abc
c = collections
c.abc = collections.abc

# importing the required libraries................................
from pptx import Presentation
from pptx.util import Inches,Pt

# defining function to give required font-style and size...........
def add_font_to_presentation(presentation, font_name):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = font_name
                        run.font.size = Pt(12)  




ppt=Presentation() 

# defining the slide layout for our slides and passing its ref. to slide_layot1....
slide_layout1=ppt.slide_layouts[6] 

#adding 2 slides to out presentation.....
s1=ppt.slides.add_slide(slide_layout1)
s2=ppt.slides.add_slide(slide_layout1)

# adding the shapes i.e text box here in our slides as txtt1 and txtt2...
left = Inches(0.5)
top = Inches(0.5)
width = Inches(9)  
height = Inches(7)
txtt1=s1.shapes.add_textbox(left,top,width,height)
txtt2=s2.shapes.add_textbox(left,top,width,height)

# passing the ref. of text frame of text boxes to box1 & box2 to access our text box content....
box1=txtt1.text_frame
box2=txtt2.text_frame

# wrapping our text to fit it in slide......
box1.word_wrap=True

# opening,reading the file sample_slide1_input.txt by passing its path and passing it text in our slide1(s1)....
with open('sample_slide1_input.txt','r') as f:
    contents = f.read()
    box1.text=contents
    
# opening,reading the file sample_slide2_input.txt by passing its path and passing it text in our slide2(s2)....
box2.word_wrap=True
with open('sample_slide2_input.txt','r') as k:
    contents=k.read()
    box2.text=contents

# Specify the font name................
font_name ='Love Ya Like A Sister'


# Apply the font to the presentation by passing it in previously defined function..........
add_font_to_presentation(ppt, font_name)

# saving file...........
ppt.save('PPT.pptx')




